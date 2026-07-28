"""產生 Ingest 測試需要的小 Raw Material。

每個格式一個 builder，全部用容器裡已有的依賴產生，不放二進位檔進版控。
"""

import json
import struct
import zipfile
import zlib


def write_pdf(path, text):
    """最小的單頁 PDF，內文用 Helvetica 印一行 ASCII 字串。"""
    content = f"BT /F1 24 Tf 72 700 Td ({text}) Tj ET".encode("ascii")
    bodies = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]"
        b" /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(content), content),
    ]

    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for number, body in enumerate(bodies, start=1):
        offsets.append(len(out))
        out += b"%d 0 obj\n%s\nendobj\n" % (number, body)

    xref_offset = len(out)
    out += b"xref\n0 %d\n" % (len(bodies) + 1)
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += b"%010d 00000 n \n" % offset
    out += b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF\n" % (
        len(bodies) + 1,
        xref_offset,
    )

    path.write_bytes(bytes(out))


def write_broken_pdf(path):
    """看起來是 PDF、但轉檔一定會失敗的檔案。"""
    path.write_bytes(b"%PDF-1.4\n" + bytes([0xFF, 0xFE, 0xFD, 0x81, 0x8F]) * 40)


def write_docx(path, text):
    from docx import Document

    document = Document()
    document.add_paragraph(text)
    document.save(str(path))


def write_pptx(path, text):
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = text
    presentation.save(str(path))


def write_xlsx(path, text):
    from openpyxl import Workbook

    workbook = Workbook()
    workbook.active["A1"] = text
    workbook.save(str(path))


def write_html(path, text):
    path.write_text(
        f"<html><body><h1>{text}</h1></body></html>",
        encoding="utf-8",
    )


def write_csv(path, text):
    path.write_text(f"姓名,單位\n{text},工程部\n", encoding="utf-8")


def write_json(path, text):
    path.write_text(json.dumps({"decision": text}, ensure_ascii=False), encoding="utf-8")


def write_xml(path, text):
    path.write_text(f"<notes><item>{text}</item></notes>", encoding="utf-8")


def write_epub(path, text):
    chapter = (
        '<?xml version="1.0" encoding="utf-8"?>'
        '<html xmlns="http://www.w3.org/1999/xhtml"><body>'
        f"<h1>{text}</h1></body></html>"
    )
    opf = (
        '<?xml version="1.0" encoding="utf-8"?>'
        '<package xmlns="http://www.idpf.org/2007/opf" version="3.0"'
        ' unique-identifier="id">'
        '<metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
        "<dc:title>會議守則</dc:title><dc:language>zh-TW</dc:language>"
        '<dc:identifier id="id">handbook</dc:identifier></metadata>'
        '<manifest><item id="c1" href="chapter1.xhtml"'
        ' media-type="application/xhtml+xml"/></manifest>'
        '<spine><itemref idref="c1"/></spine></package>'
    )
    container = (
        '<?xml version="1.0" encoding="utf-8"?>'
        '<container xmlns="urn:oasis:names:tc:opendocument:xmlns:container"'
        ' version="1.0"><rootfiles>'
        '<rootfile full-path="OEBPS/content.opf"'
        ' media-type="application/oebps-package+xml"/>'
        "</rootfiles></container>"
    )

    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("mimetype", "application/epub+zip")
        archive.writestr("META-INF/container.xml", container)
        archive.writestr("OEBPS/content.opf", opf)
        archive.writestr("OEBPS/chapter1.xhtml", chapter)


def write_png(path):
    def chunk(kind, payload):
        body = kind + payload
        return struct.pack(">I", len(payload)) + body + struct.pack(">I", zlib.crc32(body))

    header = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    pixels = zlib.compress(b"\x00\xff\xff\xff")
    path.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", header)
        + chunk(b"IDAT", pixels)
        + chunk(b"IEND", b"")
    )


# --- .msg（OLE2 複合檔）------------------------------------------------------
# markitdown 用 olefile 讀 .msg，而依賴裡沒有任何能「寫」OLE2 的套件，
# 所以這裡手工組一份最小的複合檔。每個 stream 都墊到 mini stream 門檻
# （4096 bytes）以上，就不必再實作 MiniFAT——markitdown 讀出來會 strip 掉空白。

_SECTOR = 512
_MINI_CUTOFF = 4096
_FREESECT = 0xFFFFFFFF
_ENDOFCHAIN = 0xFFFFFFFE
_FATSECT = 0xFFFFFFFD

_SUBJECT_STREAM = "__substg1.0_0037001F"
_SENDER_STREAM = "__substg1.0_0C1F001F"
_BODY_STREAM = "__substg1.0_1000001F"


def _directory_entry(name, kind, left, right, child, start, size):
    encoded = name.encode("utf-16-le") + b"\x00\x00"
    entry = (
        encoded.ljust(64, b"\x00")
        + struct.pack("<H", len(encoded))
        + bytes([kind, 1])  # 物件型別 + 顏色（black）
        + struct.pack("<III", left, right, child)
        + b"\x00" * 16  # CLSID
        + struct.pack("<I", 0)  # state bits
        + b"\x00" * 16  # 建立與修改時間
        + struct.pack("<I", start)
        + struct.pack("<Q", size)
    )
    assert len(entry) == 128
    return entry


def write_msg(path, sender, subject, body):
    # stream 名稱依 CFB 規定要照長度、大寫排序；這三個等長，字面排序即正確。
    streams = [
        (_SUBJECT_STREAM, subject),
        (_SENDER_STREAM, sender),
        (_BODY_STREAM, body),
    ]

    fat_sector_count = 1
    first_data_sector = fat_sector_count + 1  # FAT 之後是目錄，目錄之後才是資料

    payload = b""
    directory = [
        _directory_entry("Root Entry", 5, _FREESECT, _FREESECT, 1, _ENDOFCHAIN, 0)
    ]
    fat = [_FATSECT] * fat_sector_count + [_ENDOFCHAIN]

    for index, (name, text) in enumerate(streams, start=1):
        raw = text.encode("utf-16-le")
        size = max(_MINI_CUTOFF, len(raw))
        if size % _SECTOR:
            size += _SECTOR - size % _SECTOR
        padded = raw + b"\x20\x00" * ((size - len(raw)) // 2)

        start = first_data_sector + len(payload) // _SECTOR
        sectors = len(padded) // _SECTOR
        fat += list(range(start + 1, start + sectors)) + [_ENDOFCHAIN]
        payload += padded

        sibling = index + 1 if index < len(streams) else _FREESECT
        directory.append(
            _directory_entry(name, 2, _FREESECT, sibling, _FREESECT, start, len(padded))
        )

    header = (
        b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
        + b"\x00" * 16  # CLSID
        + struct.pack("<HHH", 0x003E, 0x0003, 0xFFFE)
        + struct.pack("<HH", 9, 6)  # sector shift / mini sector shift
        + b"\x00" * 6  # reserved
        + struct.pack("<I", 0)  # 目錄 sector 數：v3 固定為 0
        + struct.pack("<I", fat_sector_count)
        + struct.pack("<I", fat_sector_count - 1 + 1)  # 第一個目錄 sector
        + struct.pack("<I", 0)  # transaction signature
        + struct.pack("<I", _MINI_CUTOFF)
        + struct.pack("<I", _ENDOFCHAIN)  # 第一個 MiniFAT sector：不使用
        + struct.pack("<I", 0)
        + struct.pack("<I", _ENDOFCHAIN)  # 第一個 DIFAT sector：不使用
        + struct.pack("<I", 0)
        + struct.pack("<I", 0)  # DIFAT[0] = FAT 放在 sector 0
        + struct.pack("<I", _FREESECT) * 108
    )
    assert len(header) == _SECTOR

    fat += [_FREESECT] * (_SECTOR // 4 * fat_sector_count - len(fat))
    fat_bytes = b"".join(struct.pack("<I", entry) for entry in fat)

    directory_bytes = b"".join(directory).ljust(_SECTOR, b"\x00")

    path.write_bytes(header + fat_bytes + directory_bytes + payload)
